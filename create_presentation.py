"""
Enhanced PowerPoint presentation for GPU-Accelerated AI-Powered Healthcare Chatbot
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define consistent color scheme
PRIMARY = RGBColor(20, 184, 166)  # Teal
SECONDARY = RGBColor(59, 130, 246)  # Blue
ACCENT = RGBColor(139, 92, 246)  # Purple
DARK = RGBColor(30, 41, 59)
LIGHT_BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(71, 85, 105)

def add_background_shape(slide):
    """Add consistent background design to all slides"""
    # Top accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    # Bottom right circle decoration
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.5), Inches(6.5),
        Inches(2), Inches(2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY
    circle.fill.transparency = 0.85
    circle.line.fill.background()

def add_title_slide(title, subtitle):
    """Add an interactive title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background effect
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    
    # Decorative circles
    for i, (x, y, size, color) in enumerate([
        (8, 1, 2.5, PRIMARY),
        (1, 5, 2, SECONDARY),
        (7.5, 5.5, 1.5, ACCENT)
    ]):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.fill.transparency = 0.8
        circle.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_dict):
    """Add a content slide with icon boxes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    add_background_shape(slide)
    
    # Title with underline
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    # Underline accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(2), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()
    
    # Content
    y_position = 2.2
    for section, points in content_dict.items():
        # Section header with icon box
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y_position),
            Inches(0.4), Inches(0.4)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = PRIMARY
        icon_box.line.fill.background()
        
        # Section title
        section_box = slide.shapes.add_textbox(
            Inches(1.3), Inches(y_position),
            Inches(8), Inches(0.4)
        )
        section_frame = section_box.text_frame
        section_frame.text = section
        p = section_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK
        
        y_position += 0.5
        
        # Bullet points
        for point in points:
            point_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(y_position),
                Inches(7.5), Inches(0.35)
            )
            point_frame = point_box.text_frame
            point_frame.text = f"• {point}"
            p = point_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_GRAY
            
            y_position += 0.35
        
        y_position += 0.2

def add_two_column_slide(title, left_content, right_content):
    """Add a two-column slide with card design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    add_background_shape(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    # Left card
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.8),
        Inches(4.3), Inches(5)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = LIGHT_BG
    left_card.line.color.rgb = PRIMARY
    left_card.line.width = Pt(2)
    
    # Right card
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.8),
        Inches(4.3), Inches(5)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_BG
    right_card.line.color.rgb = SECONDARY
    right_card.line.width = Pt(2)
    
    # Add content to cards
    def add_card_content(x_pos, content, color):
        y_pos = 2.1
        for section, points in content.items():
            # Section header
            header_box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(y_pos),
                Inches(4), Inches(0.4)
            )
            header_frame = header_box.text_frame
            header_frame.text = section
            p = header_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = color
            
            y_pos += 0.45
            
            # Points
            for point in points:
                point_box = slide.shapes.add_textbox(
                    Inches(x_pos), Inches(y_pos),
                    Inches(4), Inches(0.3)
                )
                point_frame = point_box.text_frame
                point_frame.text = f"• {point}"
                p = point_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = TEXT_GRAY
                
                y_pos += 0.3
            
            y_pos += 0.15
    
    add_card_content(0.7, left_content, PRIMARY)
    add_card_content(5.4, right_content, SECONDARY)

# Slide 1: Title Slide
add_title_slide(
    "MedNexus AI",
    "GPU-Accelerated Healthcare Chatbot | Intelligent Medical Assistant"
)

# Slide 2: Project Overview
add_content_slide(
    "What is MedNexus AI?",
    {
        "Project Vision": [
            "An intelligent conversational assistant for preliminary healthcare guidance",
            "Combines advanced AI with GPU acceleration for real-time responses",
            "Bridges the gap between patients and healthcare providers"
        ],
        "Core Objectives": [
            "Provide 24/7 accessible medical information and triage",
            "Deliver instant symptom assessment with urgency classification",
            "Enable visual medical analysis through image recognition",
            "Optimize performance using GPU-accelerated AI inference"
        ]
    }
)

# Slide 3: Key Features Explained
add_content_slide(
    "Key Features & Capabilities",
    {
        "Intelligent Symptom Triage": [
            "Multi-step questioning protocol to understand symptoms deeply",
            "RED/YELLOW/GREEN urgency classification system",
            "Clinical reasoning with context-aware follow-up questions"
        ],
        "Visual Medical Analysis": [
            "Upload images for rash, injury, or symptom identification",
            "AI-powered pattern recognition and description",
            "Provides visual symptom analysis and initial assessment"
        ],
        "Medical Knowledge Base": [
            "Drug information, interactions, and side effects",
            "Treatment explanations in simple, understandable terms",
            "Health education, lifestyle tips, and preventive care"
        ]
    }
)

# Slide 4: Technology Stack
add_two_column_slide(
    "Technology Stack",
    {
        "Frontend Layer": [
            "React 19.2.1 - Modern UI framework",
            "TypeScript - Type-safe development",
            "Vite - Lightning-fast build tool",
            "Tailwind CSS - Responsive styling"
        ],
        "AI & Backend": [
            "Google Gemini 2.5 Flash Model",
            "Google GenAI SDK integration",
            "RESTful API architecture",
            "Real-time message streaming"
        ]
    },
    {
        "GPU Acceleration": [
            "CUDA/OpenCL programming",
            "Parallel processing optimization",
            "Memory bandwidth utilization",
            "Batch inference processing"
        ],
        "AI Technologies": [
            "Natural Language Processing (NLP)",
            "Deep Learning neural networks",
            "Computer Vision for images",
            "Predictive health modeling"
        ]
    }
)

# Slide 5: How GPU Acceleration Works
add_content_slide(
    "Understanding GPU Acceleration",
    {
        "What is GPU Acceleration?": [
            "GPUs have thousands of cores vs CPU's few cores",
            "Processes multiple AI calculations simultaneously (parallel processing)",
            "Dramatically reduces response time from seconds to milliseconds"
        ],
        "Benefits for Healthcare AI": [
            "Faster symptom analysis and triage decisions",
            "Support for multiple users at the same time",
            "Real-time image processing and analysis",
            "Handles complex medical queries without delays"
        ],
        "Technical Implementation": [
            "Optimized matrix operations for neural networks",
            "Efficient memory usage for large AI models",
            "Batch processing for handling multiple requests"
        ]
    }
)

# Slide 6: System Architecture
add_content_slide(
    "System Architecture Overview",
    {
        "Three-Layer Design": [
            "Presentation Layer: User-friendly React interface",
            "Application Layer: Business logic and AI integration",
            "Processing Layer: GPU-accelerated model inference"
        ],
        "Key Components": [
            "Authentication Service - Secure user login and data protection",
            "Gemini AI Service - Natural conversation and medical reasoning",
            "Message Manager - Real-time chat history and state",
            "Image Processor - Visual analysis pipeline with GPU support"
        ],
        "Data Flow": [
            "User input → AI processing → GPU acceleration → Response",
            "Images processed through specialized computer vision pipeline",
            "All data encrypted and securely transmitted"
        ]
    }
)

# Slide 7: AI Capabilities Deep Dive
add_content_slide(
    "Advanced AI Capabilities",
    {
        "Natural Language Understanding": [
            "Understands medical terminology and layman descriptions",
            "Context-aware: remembers conversation history",
            "Handles multi-turn dialogues with follow-up questions"
        ],
        "Triage Intelligence": [
            "Systematic symptom discovery through targeted questions",
            "Risk assessment based on symptom severity and combinations",
            "Emergency detection with immediate action recommendations"
        ],
        "Computer Vision": [
            "Analyzes uploaded medical images (rashes, wounds, etc.)",
            "Pattern recognition for common visual symptoms",
            "Generates detailed descriptions for healthcare providers"
        ]
    }
)

# Slide 8: Real-World Use Cases
add_two_column_slide(
    "Practical Use Cases",
    {
        "Everyday Healthcare": [
            "Pre-consultation symptom check",
            "Understanding medication instructions",
            "Checking drug interactions",
            "General health questions",
            "Appointment preparation",
            "Treatment plan explanations"
        ],
        "Preventive Care": [
            "Diet and nutrition guidance",
            "Exercise recommendations",
            "Health screening reminders",
            "Lifestyle modification tips"
        ]
    },
    {
        "Emergency Situations": [
            "Immediate symptom severity check",
            "Red flag identification",
            "Emergency service guidance",
            "First-aid instructions",
            "Hospital vs urgent care routing"
        ],
        "Chronic Condition Support": [
            "Symptom tracking assistance",
            "Medication adherence help",
            "Treatment side effect queries",
            "Lifestyle management tips"
        ]
    }
)

# Slide 9: Safety & Compliance
add_content_slide(
    "Safety, Privacy & Compliance",
    {
        "Medical Safety Protocols": [
            "Clear disclaimer: AI assistant, not a doctor replacement",
            "Always recommends professional consultation for serious issues",
            "Emergency override system for life-threatening symptoms",
            "Conservative approach to medical advice"
        ],
        "Data Privacy & Security": [
            "Secure authentication and user verification",
            "End-to-end encryption for all communications",
            "HIPAA-compliant design principles",
            "No sharing of personal health information"
        ],
        "Ethical AI Principles": [
            "Transparent about AI limitations and capabilities",
            "Educates users on appropriate use cases",
            "Bias-aware medical guidance",
            "Responsible health information delivery"
        ]
    }
)

# Slide 10: Performance Metrics
add_two_column_slide(
    "Performance & Impact",
    {
        "Speed & Efficiency": [
            "Average response: < 2 seconds",
            "GPU inference: < 500ms",
            "Image analysis: < 3 seconds",
            "Supports 100+ concurrent users"
        ],
        "Accuracy Metrics": [
            "Medical domain specialization",
            "Context-appropriate responses",
            "High emergency detection rate",
            "Consistent triage classification"
        ]
    },
    {
        "User Impact": [
            "24/7 healthcare accessibility",
            "Reduced wait times for info",
            "Better-informed patients",
            "Improved pre-consultation prep"
        ],
        "System Reliability": [
            "99.9% uptime target",
            "Automatic error recovery",
            "Load balancing for scale",
            "Performance monitoring"
        ]
    }
)

# Slide 11: Future Enhancements
add_content_slide(
    "Future Roadmap & Enhancements",
    {
        "Next-Generation Features": [
            "Voice input and speech synthesis for accessibility",
            "Multi-language support (Hindi, Spanish, Mandarin, etc.)",
            "Integration with Electronic Health Records (EHR)",
            "Telemedicine platform connectivity for seamless doctor handoff"
        ],
        "Advanced GPU Optimization": [
            "Custom CUDA kernels for faster processing",
            "Model quantization for efficiency",
            "Distributed GPU inference for global scale",
            "Edge computing deployment for remote areas"
        ],
        "Enhanced AI Capabilities": [
            "Fine-tuned models on medical datasets",
            "Advanced visual analysis (X-rays, scans)",
            "Predictive health analytics and risk scoring",
            "Personalized health recommendations based on history"
        ]
    }
)

# Slide 12: Conclusion & Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(7.5)
)
bg.fill.solid()
bg.fill.fore_color.rgb = LIGHT_BG
bg.line.fill.background()

# Decorative elements
for i, (x, y, size, color) in enumerate([
    (7.5, 0.5, 2.5, PRIMARY),
    (0.5, 5, 2, SECONDARY),
    (8, 5.5, 1.5, ACCENT)
]):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y),
        Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.fill.transparency = 0.8
    circle.line.fill.background()

# Main title
title_box = slide.shapes.add_textbox(
    Inches(1), Inches(1.5),
    Inches(8), Inches(1)
)
title_frame = title_box.text_frame
title_frame.text = "Thank You!"
p = title_frame.paragraphs[0]
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER

# Summary box
summary_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(2), Inches(3),
    Inches(6), Inches(2.5)
)
summary_box.fill.solid()
summary_box.fill.fore_color.rgb = WHITE
summary_box.line.color.rgb = PRIMARY
summary_box.line.width = Pt(3)

# Summary content
summary_text = slide.shapes.add_textbox(
    Inches(2.3), Inches(3.3),
    Inches(5.4), Inches(2)
)
text_frame = summary_text.text_frame
text_frame.text = "MedNexus AI - Making Healthcare Accessible"

# Add bullet points
for point in [
    "GPU-accelerated for real-time responses",
    "Intelligent symptom triage & visual analysis",
    "Safe, secure, and user-friendly",
    "Bridging technology and healthcare"
]:
    p = text_frame.add_paragraph()
    p.text = f"✓ {point}"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY
    p.space_before = Pt(8)

# Questions text
questions_box = slide.shapes.add_textbox(
    Inches(1), Inches(6),
    Inches(8), Inches(0.6)
)
questions_frame = questions_box.text_frame
questions_frame.text = "Questions & Discussion"
p = questions_frame.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "MedNexus_AI_Healthcare_Chatbot_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📊 Total slides: 12")
print(f"🎨 Consistent design applied across all slides")
print(f"📝 Enhanced content with clear explanations")